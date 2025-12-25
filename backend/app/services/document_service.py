"""
Document processing services
"""

import os
from pathlib import Path
from typing import Optional
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation


class DocumentService:
    """Service for document processing and text extraction"""
    
    def extract_text(self, file_path: str, file_type: str) -> str:
        """Extract text from document"""
        try:
            if file_type == "pdf":
                return self._extract_from_pdf(file_path)
            elif file_type == "docx":
                return self._extract_from_docx(file_path)
            elif file_type == "pptx":
                return self._extract_from_pptx(file_path)
            elif file_type == "txt":
                return self._extract_from_txt(file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_type}")
        except Exception as e:
            print(f"Error extracting text: {e}")
            return ""
    
    def _extract_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF"""
        text = ""
        with open(file_path, "rb") as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        return text
    
    def _extract_from_docx(self, file_path: str) -> str:
        """Extract text from DOCX"""
        doc = DocxDocument(file_path)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        return text
    
    def _extract_from_pptx(self, file_path: str) -> str:
        """Extract text from PPTX"""
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    
    def _extract_from_txt(self, file_path: str) -> str:
        """Extract text from TXT"""
        with open(file_path, "r", encoding="utf-8") as file:
            return file.read()

